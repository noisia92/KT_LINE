from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_AUTO_SIZE
from pptx.util import Inches,Pt,Cm
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
# PPT문서 생성

#캐리어 샘플
CARRIER = ['남청중-10R-LX13']
#Node와 Edge는 네오스에서 다운받은 파일에 있는 내용 그대로를 사용함
NODES = ['남청중 MSPP-L-2810-01 서청주 MSPP-L-0217-02-02 진천 MSPP-L-0607-02 괴산 MSPP-L-0215-02 증평 MSPP-L-0704-02']
EDGES = ['괴산-음성-FO-D001-020(24,24) @ 남청중-증평-FO-D003-022(48,48) @ 서청주-진천-FO-D001-020(72,72) @ 남청주-서청주-FO-005B-058(144,144) @ 음성-진천-FO-D001-012(72,72) @ 괴산-증평-FO-D004-004(72,72)']

#
Ring_Left = Inches(1.25)
RING_Width = Inches(7.5)
Ring_Height = Inches(4.0)
Ring_Top = Inches(2.0)

Node_Width = Inches(1.0)
Node_Height = Inches(0.75)

SHAPE_POSITIONS = [[],[],[],[],[],[],[],[]]
def IS_Same_Kuksa(kuksa,links_kuksa):
    if kuksa == '서청주' or kuksa == '청주중' :
        if '서청주' in links_kuksa.split('-')[:2] or '청주중' in links_kuksa.split('-')[:2]:
            return True
    elif  kuksa == '남청주' or kuksa == '남청중' :
        if '남청주' in links_kuksa.split('-')[:2] or '남청중' in links_kuksa.split('-')[:2]:
            return True
    elif kuksa in links_kuksa.split('-')[:2]:
            return True
    else :
        return False
def Make_Base_Slide(slide):
#
#최상단 제목

    left =  Inches(0)
    width = Inches(10)
    height = Inches(0.7)
    top = Inches(0)

    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, left, top, width, height
    )
    fill = shape.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(204,219,226)
    line = shape.line
    line.width = Pt(2.5)
    line.color.rgb = RGBColor(204,219,226)
    text_frame = shape.text_frame
    p = text_frame.paragraphs[0]
    run = p.add_run()
    run.text = '시외국간 대용량MSPP'
    font = run.font
    font.name = '맑은 고딕'
    font.size = Pt(24)
    font.bold = True
    font.color.rgb = RGBColor(0,0,0)

    #캐리어명 박스
    left =  Inches(0.7)
    width = Inches(3)
    height = Inches(0.5)
    top = Inches(1.0)

    shape1 = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    fill = shape1.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(204,219,226)
    line = shape1.line
    line.width = Pt(2.5)
    line.color.rgb = RGBColor(204,219,226)
    text_frame = shape1.text_frame
    p = text_frame.paragraphs[0]
    run = p.add_run()
    run.text = '▣ 캐리어명 :'
    font = run.font
    font.name = '맑은 고딕'
    font.size = Pt(18)
    font.bold = True
    font.color.rgb = RGBColor(0,0,0)

    # 링형 나타내는 모서리 둥근 사각형 큰거
    # 도형생성 - 모서리둥글고 비어있는 rectangle
    left =  Inches(1.25)
    width = Inches(7.5)
    height = Inches(4.0)
    top = Inches(2.0)

    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    fill = shape.fill
    fill.background()
    line = shape.line
    line.width = Pt(2.5)
    line.color.rgb = RGBColor(0, 163, 210)


    # 범례 나타내는  사각형
    # 도형생성 - 모서리둥글고 비어있는 rectangle
    left = Inches(1.75)
    width = Inches(6)
    height = Inches(0.5)
    top = Inches(6.4)

    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, left, top, width, height
    )
    fill = shape.fill
    fill.background()
    line = shape.line
    line.width = Pt(1.5)
    line.color.rgb = RGBColor(0, 163, 210)

    text_frame = shape.text_frame
    p = text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = '범례 |                : FDF 패치          : 광단국           : 광중계기'

    font = run.font
    font.name = '맑은 고딕'
    font.size = Pt(14)

    font.color.rgb = RGBColor(0, 0, 0)

    # 범례 도형 표기 - OVAL(타원형)
    left = Inches(3)
    width = Inches(0.3)
    height = Inches(0.2)
    top = Inches(6.55)

    shape = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, left, top, width, height
    )
    fill = shape.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(204, 219, 226)
    line = shape.line
    line.width = Pt(0.5)
    line.color.rgb = RGBColor(0, 0, 0)

    # 범례 도형 표기 - RECTANGLE
    left = Inches(4.5)
    width = Inches(0.3)
    height = Inches(0.2)
    top = Inches(6.55)

    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, left, top, width, height
    )
    fill = shape.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(204, 219, 226)
    line = shape.line
    line.width = Pt(0.5)
    line.color.rgb = RGBColor(0, 0, 0)



# 범례 도형 표기 - RIBBON
    left = Inches(6)
    width = Inches(0.4)
    height = Inches(0.2)
    top = Inches(6.55)

    shape = slide.shapes.add_shape(
        MSO_SHAPE.UP_RIBBON, left, top, width, height
    )
    fill = shape.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(204, 219, 226)
    line = shape.line
    line.width = Pt(0.5)
    line.color.rgb = RGBColor(0, 0, 0)


    # 노드 표기 - RECTANGLE1
    left = Inches(0.75)
    width = Inches(1.0)
    height = Inches(0.75)
    top = Inches(3.75)

    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, left, top, width, height
    )
    fill = shape.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(204, 219, 226)
    line = shape.line
    line.width = Pt(0.5)
    line.color.rgb = RGBColor(0, 0, 0)


    # 노드 표기 - RECTANGLE2
    left = Inches(1.125)+Ring_Left
    width = Inches(1.0)
    height = Inches(0.75)
    top = Ring_Top-(Node_Height/2)
    # top = Inches(1.625)

    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, left, top, width, height
    )
    fill = shape.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(204, 219, 226)
    line = shape.line
    line.width = Pt(0.5)
    line.color.rgb = RGBColor(0, 0, 0)


    # 노드 표기 - RECTANGLE3
    left = Inches(1.125)*2 + Ring_Left+ Node_Width
    width = Inches(1.0)
    height = Inches(0.75)
    top = Ring_Top - (Node_Height / 2)
    # top = Inches(1.625)

    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, left, top, width, height
    )
    fill = shape.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(204, 219, 226)
    line = shape.line
    line.width = Pt(0.5)
    line.color.rgb = RGBColor(0, 0, 0)

# 노드 표기 - RECTANGLE4
    left = Inches(1.125)*3 + Ring_Left + Node_Width*2
    width = Inches(1.0)
    height = Inches(0.75)
    top = Ring_Top - (Node_Height / 2)
    # top = Inches(1.625)

    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, left, top, width, height
    )
    fill = shape.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(204, 219, 226)
    line = shape.line
    line.width = Pt(0.5)
    line.color.rgb = RGBColor(0, 0, 0)


# 노드 표기 - RECTANGLE5
    left = Ring_Left + RING_Width -Node_Width/2
    width = Inches(1.0)
    height = Inches(0.75)
    top = Ring_Top + Ring_Height/2 - Node_Height/2
    # top = Inches(1.625)

    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, left, top, width, height
    )
    fill = shape.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(204, 219, 226)
    line = shape.line
    line.width = Pt(0.5)
    line.color.rgb = RGBColor(0, 0, 0)

# 노드 표기 - RECTANGLE6
    left = Inches(1.125)*3 + Ring_Left + Node_Width*2
    width = Inches(1.0)
    height = Inches(0.75)
    top = Ring_Top - (Node_Height / 2) + Ring_Height

    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, left, top, width, height
    )
    fill = shape.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(204, 219, 226)
    line = shape.line
    line.width = Pt(0.5)
    line.color.rgb = RGBColor(0, 0, 0)

# 노드 표기 - RECTANGLE7
    left = Inches(1.125)*2 + Ring_Left + Node_Width*1
    width = Inches(1.0)
    height = Inches(0.75)
    top = Ring_Top - (Node_Height / 2) + Ring_Height

    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, left, top, width, height
    )
    fill = shape.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(204, 219, 226)
    line = shape.line
    line.width = Pt(0.5)
    line.color.rgb = RGBColor(0, 0, 0)


# 노드 표기 - RECTANGLE7
    left = Inches(1.125)*1 + Ring_Left
    width = Inches(1.0)
    height = Inches(0.75)
    top = Ring_Top - (Node_Height / 2) + Ring_Height

    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, left, top, width, height
    )
    fill = shape.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(204, 219, 226)
    line = shape.line
    line.width = Pt(0.5)
    line.color.rgb = RGBColor(0, 0, 0)


#하단 조직명 나타내는 박스
    left = Inches(0)
    width = Inches(10)
    height = Inches(0.5)
    top = Inches(7)

    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, left, top, width, height
    )
    fill = shape.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(204, 219, 226)
    line = shape.line
    line.width = Pt(2.5)
    line.color.rgb = RGBColor(204, 219, 226)
    text_frame = shape.text_frame
    p = text_frame.paragraphs[0]
    run = p.add_run()
    run.text = '청주운용부'
    font = run.font
    font.name = '맑은 고딕'
    font.size = Pt(16)
    font.bold = True
    font.color.rgb = RGBColor(0, 0, 0)

# def Make_Node_Shape(Node_Info,Link_Info,Position):
#     if Position = 0 :




def Make_Base_Slide_With_Info(slide,carrier,node,edge):
    kuksa = []
    start_kuksa = ''
    second_kuksa = ''
    kuksa_links = []
    kuksa_links_sorted = []
    facility = []
    orderd_edges = []

    print(carrier)
    nodes = (node.strip().split(' '))
    edges = (edge.strip().split(' @ '))
    link_from_edge = []
    for x in edges:
        link_from_edge.append(x.strip().split('-')[0][:2]+'-'+x.strip().split('-')[1][:2])

    #중복제거하기
    node_length = (len(nodes) // 2)
    edge_length = len(edges)
    print(str(node_length) + ' ' + '개의 노드가 있습니다.')
    print(str(edge_length) + ' ' + '개의 연결구간이 있습니다.')
    if edge_length-node_length ==0:
        print('패치구간이 없습니다.')
    elif edge_length-node_length >0 :
        print(str(edge_length-node_length)+' '+'개의 패치구간이 있습니다.')
    else :
        print('구성에 오류가 있습니다.')

    kuksa_with_duplication = []
    for x in edges:
        kuksa_with_duplication.append(x.split('-')[0].strip())
        kuksa_with_duplication.append(x.split('-')[1].strip())

    print(kuksa_with_duplication)

    for i,x in enumerate(nodes) :
        if i%2==0:
            kuksa.append(x)
        else :
            facility.append(x)

    for i in range(node_length-1):
        kuksa_links.append(kuksa[i][:2]+'-'+kuksa[i+1][:2])
        temp_list = [kuksa[i][:2],kuksa[i+1][:2]]
        temp_list.sort()
        kuksa_links_sorted.append(temp_list[0][:2]+'-'+temp_list[1][:2])
    kuksa_links.append(kuksa[-1][:2] + '-' + kuksa[0][:2])
    temp_list = [kuksa[-1], kuksa[0]]
    temp_list.sort()
    kuksa_links_sorted.append(temp_list[0][:2] + '-' + temp_list[1][:2])

    print('** 광단국 순서 **')
    # print(kuksa)
    kuksa_two_character = []
    for x in kuksa:
        kuksa_two_character.append(x[:2])
    print(kuksa_two_character)
    # print('\n')
    print('** E2E 연결내역 **')
    print(edges)
    print('\n')
    # 두글자만 뽑아서 만든 E2E연결내역 국사간 링크
    print(link_from_edge)
    # 광단국만 순수하게 연결되었을 경우 링크 순서
    print(kuksa_links_sorted)
    patch_nodes = []
    for x in kuksa_links_sorted :
        if x in link_from_edge:
            pass
        else :
            print('패치구간 : ',x,'간 패치구간이 있습니다.')
            patch_nodes.append(x)

    appended_link = []
    for x in link_from_edge:
        if x in kuksa_links_sorted:
            pass
        else :
            appended_link.append(x)
            print('삽입링크',x)
    replace_tuple_list = []
    for x in patch_nodes :
        temp_link = ['', '']
        # print(kuksa_two_character.index(x.split('-')[0]))
        # print(kuksa_two_character.index(x.split('-')[1]))
        if (kuksa_two_character.index(x.split('-')[0])) < (kuksa_two_character.index(x.split('-')[1])):
            former_node = x.split('-')[0]
            latter_node = x.split('-')[1]
        else :
            former_node = x.split('-')[1]
            latter_node = x.split('-')[0]
        # print(former_node,latter_node)
        for y in appended_link:
            if (y.startswith(former_node) or y.endswith(former_node)):
                temp_link[0] = y
            elif (y.startswith(latter_node) or y.endswith(latter_node)):
                if (temp_link[1] =='') :
                    temp_link[1] = y
                else :
                    if temp_link[0].split('-')[0] in (y.split('-')) or temp_link[0].split('-')[1] in (y.split('-')) :
                        temp_link[1] = y
                    else :
                        pass

        replace_tuple_list.append((x,temp_link))

    print(replace_tuple_list)
    for x in replace_tuple_list:
        print(x[0])
        index_num = kuksa_links_sorted.index(x[0])
        print(index_num)
        kuksa_links_sorted.remove(x[0])
        kuksa_links_sorted.insert(index_num,x[1][0])
        kuksa_links_sorted.insert(index_num+1, x[1][1])
    print(kuksa_links_sorted)
    EDGES_SORTED = []

    print('정렬된 Edge')
    for x in kuksa_links_sorted:
        X_LINK1 = x.split('-')[0]
        X_Link2 = x.split('-')[1]
        for y in edges:
            Y_LINK1 = y.split('-')[0][:2]
            Y_LINK2 = y.split('-')[1][:2]
            if X_LINK1 == Y_LINK1 and X_Link2 == Y_LINK2 :
                EDGES_SORTED.append(y)
    print(EDGES_SORTED)
    left =  Inches(0)
    width = Inches(10)
    height = Inches(0.7)
    top = Inches(0)

    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, left, top, width, height
    )
    fill = shape.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(204,219,226)
    line = shape.line
    line.width = Pt(2.5)
    line.color.rgb = RGBColor(204,219,226)
    text_frame = shape.text_frame
    p = text_frame.paragraphs[0]
    run = p.add_run()
    run.text = '시외국간 대용량MSPP'
    font = run.font
    font.name = '맑은 고딕'
    font.size = Pt(24)
    font.bold = True
    font.color.rgb = RGBColor(0,0,0)

    #캐리어명 박스
    left =  Inches(0.7)
    width = Inches(5)
    height = Inches(0.5)
    top = Inches(1.0)

    shape1 = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    fill = shape1.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(204,219,226)
    line = shape1.line
    line.width = Pt(2.5)
    line.color.rgb = RGBColor(204,219,226)
    text_frame = shape1.text_frame
    text_frame.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
    p = text_frame.paragraphs[0]
    run = p.add_run()
    run.text = '▣ 캐리어명 :'+carrier
    font = run.font
    font.name = '맑은 고딕'
    font.size = Pt(16)
    font.bold = True
    font.color.rgb = RGBColor(0,0,0)

    # 링형 나타내는 모서리 둥근 사각형 큰거
    # 도형생성 - 모서리둥글고 비어있는 rectangle
    left = Inches(1.25)
    width = Inches(7.5)
    height = Inches(4.0)
    top = Inches(2.0)

    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    fill = shape.fill
    fill.background()
    line = shape.line
    line.width = Pt(2.5)
    line.color.rgb = RGBColor(0, 163, 210)

    # 범례 나타내는  사각형
    # 도형생성 - 모서리둥글고 비어있는 rectangle
    left = Inches(1.75)
    width = Inches(6)
    height = Inches(0.5)
    top = Inches(6.4)

    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, left, top, width, height
    )
    fill = shape.fill
    fill.background()
    line = shape.line
    line.width = Pt(1.5)
    line.color.rgb = RGBColor(0, 163, 210)

    text_frame = shape.text_frame
    p = text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = '범례 |                : FDF 패치          : 광단국           : 광중계기'

    font = run.font
    font.name = '맑은 고딕'
    font.size = Pt(14)

    font.color.rgb = RGBColor(0, 0, 0)

    # 범례 도형 표기 - OVAL
    left = Inches(3)
    width = Inches(0.3)
    height = Inches(0.2)
    top = Inches(6.55)

    shape = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, left, top, width, height
    )
    fill = shape.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(204, 219, 226)
    line = shape.line
    line.width = Pt(0.5)
    line.color.rgb = RGBColor(0, 0, 0)

    # 범례 도형 표기 - RECTANGLE
    left = Inches(4.5)
    width = Inches(0.3)
    height = Inches(0.2)
    top = Inches(6.55)

    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, left, top, width, height
    )
    fill = shape.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(204, 219, 226)
    line = shape.line
    line.width = Pt(0.5)
    line.color.rgb = RGBColor(0, 0, 0)

    # 범례 도형 표기 - RIBBON
    left = Inches(6)
    width = Inches(0.4)
    height = Inches(0.2)
    top = Inches(6.55)

    shape = slide.shapes.add_shape(
        MSO_SHAPE.UP_RIBBON, left, top, width, height
    )
    fill = shape.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(204, 219, 226)
    line = shape.line
    line.width = Pt(0.5)
    line.color.rgb = RGBColor(0, 0, 0)



    #
    #첫 장 광단국 한개 샘플 삽입
    #
    left = Inches(0.75)
    width = Inches(1.0)
    height = Inches(0.75)
    top = Inches(3.75)

    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, left, top, width, height
    )
    fill = shape.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(204, 219, 226)
    line = shape.line
    line.width = Pt(0.5)
    line.color.rgb = RGBColor(0, 0, 0)

    text_frame = shape.text_frame
    p = text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = nodes[0] + nodes[1]

    # 첫번째 노드의 Edge 텍스트상자
    font = run.font
    font.name = '맑은 고딕'
    font.size = Pt(11)
    font.color.rgb = RGBColor(0, 0, 0)

    txBox = slide.shapes.add_textbox(left+Inches(0.5), top-Inches(0.8), width+Inches(2), height-Inches(0.5))

    tf = txBox.text_frame
    p = tf.add_paragraph()
    p.text = edges[0]
    p.font.bold = True
    p.font.size = Pt(11)


    #하단 조직명 나타내는 박스
    left = Inches(0)
    width = Inches(10)
    height = Inches(0.5)
    top = Inches(7)

    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, left, top, width, height
    )
    fill = shape.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(204, 219, 226)
    line = shape.line
    line.width = Pt(2.5)
    line.color.rgb = RGBColor(204, 219, 226)
    text_frame = shape.text_frame
    p = text_frame.paragraphs[0]
    run = p.add_run()
    run.text = '청주운용부'
    font = run.font
    font.name = '맑은 고딕'
    font.size = Pt(16)
    font.bold = True
    font.color.rgb = RGBColor(0, 0, 0)

prs = Presentation()
default_number = 10
title_only_slide_layout = prs.slide_layouts[6]

# 슬라이드 추가
slide_list = []
for i in range(default_number) :
    slide_list.append(prs.slides.add_slide(title_only_slide_layout))

for i,x in enumerate(slide_list):
    if i ==0:
        Make_Base_Slide_With_Info(x,CARRIER[0],NODES[0],EDGES[0])
    else :
        Make_Base_Slide(x)
prs.save('test.pptx')